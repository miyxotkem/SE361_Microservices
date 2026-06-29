from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_slide(prs, title, content_bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
            
        p.text = bullet
        p.font.size = Pt(20)
        
    return slide

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

prs = Presentation()

# Slide 1
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Hệ Thống E-Learning SmartEdu"
subtitle.text = "Kiến trúc Microservices\n(Báo cáo Đồ án Cuối kỳ)"
add_notes(slide1, "Chào thầy và các bạn. Hôm nay nhóm chúng em xin trình bày về SmartEdu - một hệ thống E-learning. Thay vì làm một trang web truyền thống nơi mọi tính năng nhồi nhét vào một khối duy nhất, nhóm em quyết định áp dụng Kiến trúc Microservices. Giống như một nhà hàng: thay vì 1 người làm mọi việc, em chia ra thành các chuyên gia riêng biệt: đầu bếp, phục vụ, bảo vệ... Mỗi dịch vụ làm một việc tốt nhất có thể.")

# Slide 2
s2 = create_slide(prs, "Kiến trúc Tổng thể & YARP API Gateway", [
    "• Cửa ngõ (API Gateway): YARP (Yet Another Reverse Proxy - Port 7000)",
    "• 6 Services độc lập: Identity, Course, Exam, Notification, Comment, Payment",
    "• Tránh gọi trực tiếp: Client (WPF Desktop App) chỉ giao tiếp qua Gateway",
    "• Dễ dàng phân tải, bảo mật và điều hướng (Routing)."
])
add_notes(s2, "Đây là bức tranh toàn cảnh. Mọi yêu cầu từ người dùng (Client) không bao giờ đi thẳng tới các dịch vụ. Nó phải qua một 'Cô tiếp tân' tên là API Gateway. Ví dụ: Khi xem khóa học, gọi Gateway/api/courses. Cô tiếp tân sẽ điều hướng yêu cầu đến phòng ban Course. Giúp bảo mật thông tin nội bộ và phân tải khi đông khách.")

# Slide 3
s3 = create_slide(prs, "Dữ liệu Độc lập (Database-per-Service & Polyglot Persistence)", [
    "• Tuân thủ nghiêm ngặt nguyên tắc: Mỗi service quản lý 1 Database riêng.",
    "• Relational DB (PostgreSQL): Dùng cho Identity và Payment (Cần ACID chặt chẽ, schema tĩnh).",
    "• NoSQL (Google Firestore): Dùng cho Course, Exam, Notification, Comment (Cấu trúc linh hoạt).",
    "• Cache/State Store (Redis): Quản lý trạng thái luồng đăng ký học."
])
add_notes(s3, "Nhóm em tuân thủ 100% Database Isolation. Dữ liệu về Tiền bạc (Payment) và Tài khoản (Identity) cần chính xác tuyệt đối -> dùng PostgreSQL (như một tủ hồ sơ kế toán). Dữ liệu Bài giảng (Course) linh hoạt định dạng -> dùng Firestore (như một ngăn kéo không giới hạn định dạng). Điều này gọi là Polyglot Persistence.")

# Slide 4
s4 = create_slide(prs, "Giao tiếp Đồng bộ (Synchronous - gRPC)", [
    "• Vấn đề: Database bị cô lập, truy vấn chéo rất tốn thời gian nếu dùng REST API HTTP/1.1.",
    "• Giải pháp: Dùng gRPC (Remote Procedure Call) qua giao thức HTTP/2.",
    "• Ví dụ: Exam Service gọi gRPC tới Course Service để lấy thông tin khóa học/sinh viên nhanh chóng.",
    "• Ưu điểm: Tốc độ cực nhanh, payload nhị phân nhẹ."
])
add_notes(s4, "Khi 'Dịch vụ Thi' muốn biết sinh viên có thuộc 'Dịch vụ Khóa học' không, làm sao lấy dữ liệu khi DB bị khóa kín? Nhóm em dùng gRPC. Nó giống như đường dây điện thoại hotline mã hóa tốc độ cao nối trực tiếp giữa hai phòng ban. Dữ liệu truyền đi dưới dạng nhị phân, cực kỳ nhẹ và nhanh, không bị nghẽn mạng.")

# Slide 5
s5 = create_slide(prs, "Giao tiếp Bất đồng bộ (Event-Driven với RabbitMQ)", [
    "• Message Broker: Sử dụng RabbitMQ.",
    "• Pattern: Publish/Subscribe (Pub/Sub).",
    "• Event: Ví dụ ExamPublishedEvent, PaymentCompletedEvent.",
    "• Ưu điểm: Các dịch vụ không bị phụ thuộc vòng (Decoupled), không cần đợi nhau."
])
add_notes(s5, "Với các tác vụ không cần phản hồi ngay, em dùng kiến trúc Hướng sự kiện (RabbitMQ). Ví dụ: Khi Giảng viên 'Công bố Đề thi'. Dịch vụ Exam chỉ cầm loa hét lên (Publish Event) vào RabbitMQ: 'Có đề thi mới!'. Dịch vụ Notification luôn nghe loa, khi nghe thấy sẽ tự động đi gửi Email cho sinh viên. Giúp các dịch vụ Decoupled (Tách rời).")

# Slide 6
s6 = create_slide(prs, "Bài toán Giao dịch Phân tán (Distributed Transactions)", [
    "• Vấn đề lớn nhất của Microservices: Mua khóa học gồm 2 bước ở 2 DB.",
    "• Bước 1: Thanh toán thành công (Lưu ở Payment DB).",
    "• Bước 2: Cấp quyền học (Lưu ở Course DB).",
    "• Rủi ro: Trả tiền xong nhưng Course DB sập thì sao? Không thể Rollback dễ dàng như hệ thống nguyên khối (Monolith)."
])
add_notes(s6, "Khi Mua khóa học, có 2 phần: Thanh toán (DB Payment) và Cấp quyền (DB Course). Ví dụ sinh viên quẹt VNPay 500k xong, nhưng Dịch vụ Course bị sập, không mở khóa học. Sinh viên mất tiền mà không được học. Ta không thể gọi lệnh Rollback DB của Payment từ Course được! Đó là bài toán Giao dịch phân tán.")

# Slide 7
s7 = create_slide(prs, "Giải pháp: Saga Pattern & Hoàn tiền tự động", [
    "• Áp dụng mô hình: Saga Orchestration (Sử dụng MassTransit StateMachine).",
    "• Trình quản lý (Orchestrator): Đặt tại Course Service, lưu trạng thái bằng Redis.",
    "• Bù đắp (Compensating Transaction): Nếu Course không thể cấp quyền, nó tự động bắn lệnh `RefundPaymentCommand` về Payment.",
    "• Payment lắng nghe lệnh và hoàn lại tiền cho sinh viên. Tính nhất quán được bảo đảm!"
])
add_notes(s7, "Em cài đặt Saga Pattern. Course đóng vai trò 'Người nhạc trưởng' (Orchestrator). Khi nhận tin 'Đã trả tiền', Nhạc trưởng ra lệnh Mở khóa học. Nếu Mở khóa học lỗi, Nhạc trưởng lập tức thực hiện Compensating Transaction: bắn lệnh Hoàn tiền (Refund) về Payment Service yêu cầu trả lại 500k. Nhờ vậy, dữ liệu toàn hệ thống luôn nhất quán, không ai bị mất tiền oan!")

# Slide 8
s8 = create_slide(prs, "Kiến trúc Code: Clean Architecture tại Exam.API", [
    "• Thay vì rập khuôn, Exam.API (phức tạp) được thiết kế theo Clean Architecture.",
    "• Bao gồm 4 Layers: API, Application, Domain, Infrastructure.",
    "• Dependency Inversion: Lõi Domain không phụ thuộc vào bất cứ công nghệ nào.",
    "• Dễ dàng bảo trì và thay đổi Database (Infrastructure) trong tương lai."
])
add_notes(s8, "Tầng Domain (Lõi) chứa các quy tắc nghiệp vụ (VD: Điểm thi không được < 0). Tầng này không hề biết Database là PostgreSQL hay Firestore. Nếu sếp yêu cầu đổi Database, em chỉ cần sửa tầng vỏ (Infrastructure) mà không đụng một dòng logic nào bên trong.")

# Slide 9
s9 = create_slide(prs, "Độ tin cậy: Exception Handling & Health Checks", [
    "• Bắt lỗi tập trung (Global Exception Handler): Bắt mọi lỗi và trả về chuẩn ProblemDetails (RFC 7807).",
    "• Không lộ lỗi nhạy cảm ra ngoài, giúp Client dễ dàng parse lỗi.",
    "• Khám sức khỏe định kỳ (Health Checks): Ping liên tục tới PostgreSQL, Firestore, Redis mỗi 30s.",
    "• Cập nhật trạng thái hệ thống qua Docker Compose."
])
add_notes(s9, "Để tránh việc app crash báo lỗi khó hiểu. Em xây dựng Global Exception Handler, mọi lỗi (sai pass, đứt mạng) đều bị tóm lại và định dạng đẹp mắt trước khi báo cho người dùng. Các Health Checks tự động ping DB mỗi 30s, nếu chết nó sẽ giơ cờ báo động ngay lập tức.")

# Slide 10
s10 = create_slide(prs, "Nhìn nhận Hạn chế & Nợ kỹ thuật (Technical Debt)", [
    "• Bộ nhớ đệm Thanh toán: Payment API đang dùng IMemoryCache (Nguy cơ mất map giao dịch khi restart -> Cần chuyển sang Redis).",
    "• Bảo mật (Security): Payment API thiếu màng lọc JWT Auth (Lỗ hổng bảo mật cần vá ngay).",
    "• Quản lý thông tin mật: Các secret key của DB, Firebase, VNPay đang hardcode trong cấu hình, cần đưa vào Secret Manager."
])
add_notes(s10, "Nhóm em nhận thức hệ thống còn các 'Nợ kỹ thuật'. Thứ nhất, luồng thanh toán VNPay đang lưu trạng thái tạm vào RAM (IMemoryCache), nếu Server khởi động lại, dữ liệu sẽ mất -> Cần chuyển sang Redis. Thứ hai, Payment Service chưa tích hợp bảo mật Token, các mật khẩu đang bị hardcode. Trong phiên bản tới, chúng em sẽ đưa lên Secret Manager.")

output_path = os.path.join(os.getcwd(), 'SmartEdu_Microservices_Presentation.pptx')
prs.save(output_path)
print(f"Slides saved to {output_path}")
