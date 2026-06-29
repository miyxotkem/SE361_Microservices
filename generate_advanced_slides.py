from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_slide(prs, title, content_bullets):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
            
        p.text = bullet
        p.font.size = Pt(20)
        
    return slide

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

prs = Presentation()

# Slide 1
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Hệ thống Học trực tuyến SmartEdu"
subtitle.text = "Tiếp cận theo Kiến trúc Microservices\n(Báo cáo Kỹ thuật Chuyên sâu)"
add_notes(slide1, "Chào thầy và các bạn. Hôm nay em xin trình bày về kiến trúc backend của hệ thống SmartEdu. Thay vì chọn kiến trúc Monolith truyền thống dễ triển khai, nhóm quyết định tiếp cận bằng kiến trúc Microservices để giải quyết các bài toán về khả năng mở rộng (scalability), cô lập lỗi (fault isolation) và ứng dụng đa ngôn ngữ lưu trữ (polyglot persistence).")

# Slide 2
s2 = create_slide(prs, "Kiến trúc Tổng thể (High-level Architecture)", [
    "• Lớp Client: WPF Desktop App.",
    "• Lớp Gateway: YARP (Yet Another Reverse Proxy).",
    "• Lớp Services: 6 Microservices (Identity, Course, Exam, Notification, Comment, Payment).",
    "• Lớp Infrastructure: RabbitMQ (Event Bus), Redis (Saga State Store)."
])
add_notes(s2, "Hệ thống bao gồm 6 domain services giao tiếp với bên ngoài thông qua YARP Gateway. Mọi request từ Desktop App đều đi qua Gateway để được routing. Ở bên dưới, các services hoạt động hoàn toàn độc lập với nhau, giao tiếp thông qua gRPC cho các truy vấn đồng bộ và RabbitMQ cho các sự kiện bất đồng bộ.")

# Slide 3
s3 = create_slide(prs, "Điểm sáng 1: 100% Database Isolation & Polyglot Persistence", [
    "• 100% Database Isolation: Không có chuyện 2 service gọi chung 1 database.",
    "• PostgreSQL (Supabase): Cho Identity và Payment (Cần tính ACID, quan hệ chặt chẽ).",
    "• Firestore (NoSQL): Cho Course, Exam, Comment, Notification (Dữ liệu document, read-heavy).",
    "• Redis: Lưu trữ trạng thái của Saga State Machine (Tốc độ cao)."
])
add_notes(s3, "Điểm sáng lớn nhất của dự án là việc tuân thủ tuyệt đối quy tắc 'Database-per-Service'. Rất nhiều dự án sinh viên tự xưng là Microservices nhưng lại dùng chung một DB vật lý. Ở đây, em tách biệt hoàn toàn. Hơn thế nữa, em áp dụng 'Polyglot Persistence' – chọn DB dựa trên đặc thù dữ liệu: PostgreSQL cho giao dịch thanh toán cần ACID, và NoSQL Firestore cho dữ liệu bài giảng/đề thi dạng document.")

# Slide 4
s4 = create_slide(prs, "Điểm sáng 2: Kiến trúc Phần mềm Linh hoạt", [
    "• Exam Service: Dùng Clean Architecture (4 layers) cho domain phức tạp.",
    "• Các Service khác: Dùng CQRS Pattern kết hợp MediatR và Carter Minimal APIs.",
    "• Lợi ích: Tối ưu hiệu năng, giảm overhead, không gò bó 1 pattern cho mọi bài toán."
])
add_notes(s4, "Thay vì ép buộc một khuôn mẫu cho toàn hệ thống, em áp dụng linh hoạt các pattern. Với Exam Service – nơi chứa logic chấm điểm và sinh đề phức tạp, em dùng Clean Architecture để đảm bảo tính dễ bảo trì. Với các service như Notification hay Comment, em dùng CQRS và Minimal APIs để tối ưu hóa hiệu năng.")

# Slide 5
s5 = create_slide(prs, "Giao tiếp Liên dịch vụ (Inter-Service Communication)", [
    "• Synchronous (Đồng bộ): Sử dụng gRPC (Protobuf) cho tốc độ cao và contract rành mạch.",
    "• Asynchronous (Bất đồng bộ): Sử dụng RabbitMQ + MassTransit cho kiến trúc Event-Driven.",
    "• Real-time: Sử dụng SignalR bắn trực tiếp trạng thái về client.",
    "• Loại bỏ hoàn toàn HTTP REST calls giữa các services."
])
add_notes(s5, "Để tránh tình trạng 'Spaghetti Network', em loại bỏ hoàn toàn việc gọi REST API giữa các service. Truy vấn đồng bộ sẽ dùng gRPC vì nó dùng binary payload (Protobuf) nhanh hơn JSON nhiều lần. Các tác vụ không yêu cầu trả về ngay lập tức sẽ được publish qua RabbitMQ Message Bus.")

# Slide 6
s6 = create_slide(prs, "Giao dịch Phân tán với Saga Pattern", [
    "• Bài toán: Mua khóa học liên quan đến 2 DB khác nhau (Payment Postgres và Course Firestore).",
    "• Vấn đề: Tránh tình trạng đã trừ tiền nhưng chưa mở khóa học.",
    "• Giải pháp: Triển khai Saga Pattern (Orchestration) quản lý bởi MassTransit StateMachine."
])
add_notes(s6, "Đây là bài toán kinh điển trong Microservices: Distributed Transactions. Do Payment và Course nằm ở 2 DB khác nhau, ta không thể dùng transaction thông thường. Em đã triển khai Saga Pattern theo mô hình Orchestrator tại Course Service để quản lý toàn bộ luồng đăng ký này.")

# Slide 7
s7 = create_slide(prs, "Cơ chế Bù trừ của Saga (Compensating Transaction)", [
    "• Happy Path: Course Saga nhận PaymentCompletedEvent -> Cập nhật Firestore -> Kích hoạt Notification & SignalR.",
    "• Failure Path: Nếu cập nhật Course DB thất bại -> Saga kích hoạt Giao dịch bù trừ (Compensating Transaction).",
    "• Bắn lệnh RefundPaymentCommand về Payment API để tự động hoàn tiền cho người dùng."
])
add_notes(s7, "Giáo sư có thể hỏi: 'Nếu người dùng trả tiền xong, nhưng Course DB bị sập thì sao?'. Câu trả lời là Saga sẽ ghi nhận lỗi và tự động thực thi 'Compensating Transaction' (Giao dịch bù trừ). Cụ thể, Course Saga sẽ bắn event Refund về cho Payment API để hoàn tiền cho người dùng. Trạng thái pending được lưu an toàn trong Redis.")

# Slide 8
s8 = create_slide(prs, "Xử lý Lỗi & Khả năng Quan sát (Observability)", [
    "• Global Exception Handler: Trả lỗi theo chuẩn HTTP RFC 7807 (ProblemDetails).",
    "• MediatR Pipeline: Tự động log và cảnh báo nếu request chạy chậm hơn 3 giây.",
    "• Health Checks: Endpoint /health theo dõi kết nối Postgres, Firestore và Redis liên tục."
])
add_notes(s8, "Hệ thống được thiết kế với cơ chế Global Exception Handler tuân theo chuẩn RFC 7807, giúp format lỗi thống nhất. Ngoài ra, pipeline chặn mọi request đi qua để đo thời gian phản hồi, tự động cảnh báo nếu tốn quá 3 giây. Health Checks đảm bảo hệ thống luôn theo dõi được trạng thái sống còn của các DB.")

# Slide 9
s9 = create_slide(prs, "Nhìn nhận Thực tế - Trade-offs & Technical Debt", [
    "• Trade-off: Polyglot Persistence mang lại tốc độ nhưng làm tăng độ phức tạp vận hành.",
    "• Technical Debt 1: Thông tin nhạy cảm (API Keys, DB Passwords) đang hardcode ở appsettings.json thay vì dùng Secret Manager.",
    "• Technical Debt 2: Payment API thiếu xác thực JWT, tiềm ẩn rủi ro bảo mật nghiêm trọng.",
    "• Technical Debt 3: Payment Cache đang dùng IMemoryCache thay vì Redis."
])
add_notes(s9, "Không có kiến trúc nào là hoàn hảo. Việc áp dụng quá nhiều công nghệ lưu trữ khiến chi phí vận hành tăng cao. Dưới góc độ của một Architect, em nhận thức rõ hệ thống vẫn còn nợ kỹ thuật: Việc chưa dùng Secret Manager và thiếu xác thực tại Payment API là những rủi ro bảo mật cần ưu tiên giải quyết. Việc tự chỉ ra những khuyết điểm này chứng tỏ mình hoàn toàn làm chủ hệ thống.")

# Slide 10
s10 = create_slide(prs, "Tổng kết", [
    "• Kiến trúc Microservices giải quyết tốt bài toán phân tải và cô lập.",
    "• Giao tiếp hiệu quả với gRPC và RabbitMQ.",
    "• Đảm bảo dữ liệu với Saga Pattern.",
    "• Cảm ơn thầy và các bạn đã lắng nghe!"
])
add_notes(s10, "Tóm lại, việc áp dụng Microservices đã giúp hệ thống đạt được tính cô lập, linh hoạt và chịu tải tốt. Dù còn một số nợ kỹ thuật cần xử lý ở phase sau, core architecture đã khá vững chắc. Em xin kết thúc phần trình bày. Kính mời thầy đặt câu hỏi.")

output_path = os.path.join(os.getcwd(), 'SmartEdu_Advanced_Presentation.pptx')
prs.save(output_path)
print(f"Slides saved to {output_path}")
